#!/usr/bin/env python3
"""Rebuild prez_horizons.pptx as a designed, visually coherent deck.

Content preserved from the original 4-slide deck (INCENTIVE Horizons contribution,
3 work packages). Restyled into a title slide, a work-package overview, and one
3-column slide each for Research questions, Methods and Impacts, plus a closing.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

EMU = 914400
def IN(x): return Emu(int(x * EMU))

# ── Palette ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0B, 0x25, 0x45)   # deep navy (titles)
INK    = RGBColor(0x22, 0x2B, 0x33)   # body text
MUTED  = RGBColor(0x6B, 0x77, 0x85)   # secondary text
LIGHT  = RGBColor(0xF4, 0xF6, 0xF8)   # slide background
CARD   = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xE2, 0xE7, 0xEC)   # card borders
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEAL   = RGBColor(0x1B, 0x99, 0x8B)   # WP1
CORAL  = RGBColor(0xE0, 0x7A, 0x5F)   # WP2
SLATE  = RGBColor(0x3D, 0x5A, 0x80)   # WP3
TEALD  = RGBColor(0x14, 0x6F, 0x65)
PILL   = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "Segoe UI Light"
HEAD_FONT  = "Segoe UI Semibold"
BODY_FONT  = "Segoe UI"

WPS = [
    {"n": "1", "title": "Public attitudes",          "meta": "24-month postdoc", "c": TEAL},
    {"n": "2", "title": "Policymakers’ thoughts", "meta": "36-month PhD",     "c": CORAL},
    {"n": "3", "title": "Integrating world models",   "meta": "18-month postdoc", "c": SLATE},
]

prs = Presentation()
prs.slide_width  = IN(13.333)
prs.slide_height = IN(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

# ── primitives ─────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, shp, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    s = slide.shapes.add_shape(shp, IN(x), IN(y), IN(w), IN(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(line_w)
    if radius is not None and shp == MSO_SHAPE.ROUNDED_RECTANGLE:
        try: s.adjustments[0] = radius
        except Exception: pass
    return s

def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is a list of (txt, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space_after is not None: p.space_after = Pt(space_after)
        if line_spacing is not None: p.line_spacing = line_spacing
        for (txt, size, color, bold, font, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.name = font; r.font.color.rgb = color
    return tb

def R(txt, size, color, bold=False, font=BODY_FONT, italic=False):
    return (txt, size, color, bold, font, italic)

def set_bullet(p, char="•", color=TEAL, font="Arial", marL=178000, indent=-178000):
    pPr = p._p.get_or_add_pPr()
    pPr.set('marL', str(marL)); pPr.set('indent', str(indent))
    for tag in ('a:buNone', 'a:buChar', 'a:buAutoNum', 'a:buFont', 'a:buClr'):
        for e in pPr.findall(qn(tag)): pPr.remove(e)
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgb = pPr.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (color[0], color[1], color[2])})
    buClr.append(srgb); pPr.append(buClr)
    pPr.append(pPr.makeelement(qn('a:buFont'), {'typeface': font}))
    pPr.append(pPr.makeelement(qn('a:buChar'), {'char': char}))

def bullets(slide, x, y, w, h, items, size=11.5, color=INK, bullet_color=TEAL):
    tb = slide.shapes.add_textbox(IN(x), IN(y), IN(w), IN(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0); tf.margin_top = tf.margin_bottom = Pt(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(7); p.line_spacing = 1.03
        r = p.add_run(); r.text = it
        r.font.size = Pt(size); r.font.name = BODY_FONT; r.font.color.rgb = color
        set_bullet(p, color=bullet_color)
    return tb

def footer(slide, num):
    box(slide, MSO_SHAPE.RECTANGLE, 0, SH - 0.12, SW, 0.12, fill=NAVY)
    for i, wp in enumerate(WPS):  # three colour ticks on the bottom band
        box(slide, MSO_SHAPE.RECTANGLE, 0 + i * (SW / 3), SH - 0.12, SW / 3, 0.12, fill=wp["c"])
    text(slide, 0.55, SH - 0.62, 8, 0.3,
         [[R("INCENTIVE Horizons  ·  Adrien Fabre (CNRS–CIRED)", 9, MUTED)]])
    text(slide, SW - 1.55, SH - 0.62, 1.0, 0.3,
         [[R(str(num), 9, MUTED, bold=True)]], align=PP_ALIGN.RIGHT)

def slide_title(slide, title):
    text(slide, 0.55, 0.42, 12, 0.8, [[R(title, 28, NAVY, bold=True, font=HEAD_FONT)]])
    box(slide, MSO_SHAPE.RECTANGLE, 0.57, 1.16, 0.62, 0.10, fill=TEAL)
    box(slide, MSO_SHAPE.RECTANGLE, 1.19, 1.16, 0.40, 0.10, fill=CORAL)
    box(slide, MSO_SHAPE.RECTANGLE, 1.59, 1.16, 0.22, 0.10, fill=SLATE)

# ── columns layout ──────────────────────────────────────────────────────────────
MARG = 0.55
GAP = 0.34
COLW = (SW - 2 * MARG - 2 * GAP) / 3
def col_x(i): return MARG + i * (COLW + GAP)

def wp_column(slide, i, y0, h, body_fn, header_h=1.02):
    wp = WPS[i]; x = col_x(i)
    # card
    box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, COLW, h, fill=CARD, line=LINE, line_w=1.0, radius=0.045)
    # colour header
    box(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, COLW, header_h, fill=wp["c"], radius=0.06)
    box(slide, MSO_SHAPE.RECTANGLE, x, y0 + header_h - 0.22, COLW, 0.22, fill=wp["c"])
    # number badge
    bd = 0.56
    box(slide, MSO_SHAPE.OVAL, x + 0.20, y0 + (header_h - bd) / 2, bd, bd, fill=WHITE)
    text(slide, x + 0.20, y0 + (header_h - bd) / 2, bd, bd, [[R(wp["n"], 22, wp["c"], bold=True, font=HEAD_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    # title + meta
    text(slide, x + 0.92, y0 + 0.12, COLW - 1.05, header_h - 0.2,
         [[R(wp["title"], 15.5, WHITE, bold=True, font=HEAD_FONT)],
          [R(wp["meta"], 10.5, WHITE)]], space_after=2, anchor=MSO_ANCHOR.MIDDLE)
    body_fn(x, y0 + header_h + 0.12)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s, WHITE)
# left ribbon (three WP colours)
seg = SH / 3
for i, wp in enumerate(WPS):
    box(s, MSO_SHAPE.RECTANGLE, 0, i * seg, 0.34, seg, fill=wp["c"])
# faint oversized monogram circles (decoration, right side)
box(s, MSO_SHAPE.OVAL, 10.7, -1.2, 3.6, 3.6, fill=RGBColor(0xEC, 0xF3, 0xF2))
box(s, MSO_SHAPE.OVAL, 11.7, 4.7, 3.0, 3.0, fill=RGBColor(0xF6, 0xEE, 0xEA))

text(s, 1.05, 1.55, 10.5, 0.5,
     [[R("CONTRIBUTION TO THE INCENTIVE HORIZONS PROJECT", 13.5, TEALD, bold=True, font=HEAD_FONT)]])
text(s, 1.0, 2.05, 10.7, 2.3,
     [[R("Understanding public and elite", 40, NAVY, bold=True, font=TITLE_FONT)],
      [R("attitudes to tax clubs", 40, NAVY, bold=True, font=TITLE_FONT)]],
     line_spacing=1.02, space_after=0)
box(s, MSO_SHAPE.RECTANGLE, 1.07, 4.32, 2.4, 0.07, fill=TEAL)
text(s, 1.05, 4.6, 10, 0.5, [[R("Adrien Fabre", 20, INK, bold=True, font=HEAD_FONT)]])
text(s, 1.05, 5.06, 10, 0.4, [[R("CNRS  ·  CIRED", 14, MUTED)]])
text(s, 1.05, 6.35, 11, 0.4,
     [[R("Three work packages:  ", 13, MUTED),
       R("public attitudes", 13, TEAL, bold=True),
       R("  ·  ", 13, MUTED), R("policymakers", 13, CORAL, bold=True),
       R("  ·  ", 13, MUTED), R("world models", 13, SLATE, bold=True)]])

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Work-package overview
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s, LIGHT)
slide_title(s, "The project in three work packages")
text(s, 0.55, 1.34, 12, 0.4,
     [[R("Public opinion, elite engagement and quantitative modelling of international tax clubs.", 13, MUTED, italic=True)]])
ov = [
    ("Representative 25k-respondent survey across 12 diverse countries, conjoint analysis on future scenarios, and simulation of respondents’ own outcomes.",
     "25k respondents · 12 countries"),
    ("One-to-one engagement with ministers’ advisors, diplomats, MPs & MEPs — questionnaires, interviews, and comparison with the public.",
     "Advisors · diplomats · MPs · MEPs"),
    ("A country-decile model of tax clubs (CGE + distributional model + tax & climate modules) and an online labour-footprint simulator.",
     "Country×decile granularity"),
]
y0, ch = 2.05, 4.55
for i, wp in enumerate(WPS):
    x = col_x(i)
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, COLW, ch, fill=CARD, line=LINE, line_w=1.0, radius=0.045)
    box(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, y0, COLW, 0.16, fill=wp["c"], radius=0.5)
    bd = 0.7
    box(s, MSO_SHAPE.OVAL, x + 0.32, y0 + 0.34, bd, bd, fill=wp["c"])
    text(s, x + 0.32, y0 + 0.34, bd, bd, [[R(wp["n"], 26, WHITE, bold=True, font=HEAD_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + 0.30, y0 + 1.18, COLW - 0.6, 0.8, [[R(wp["title"], 18, NAVY, bold=True, font=HEAD_FONT)]])
    # duration pill
    pill = box(s, MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.30, y0 + 1.78, 2.1, 0.34, fill=wp["c"], radius=0.5)
    text(s, x + 0.30, y0 + 1.78, 2.1, 0.34, [[R(wp["meta"].upper(), 9.5, WHITE, bold=True, font=HEAD_FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, x + 0.30, y0 + 2.35, COLW - 0.6, 1.5, [[R(ov[i][0], 12, INK)]], line_spacing=1.1)
    text(s, x + 0.30, y0 + ch - 0.55, COLW - 0.6, 0.4, [[R(ov[i][1], 10.5, wp["c"], bold=True, font=HEAD_FONT)]])
footer(s, 2)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Research questions
# ════════════════════════════════════════════════════════════════════════════════
RQ = [
    ["Would the public accept tax-club proposals?",
     "How would people allocate tax revenues across funds?",
     "Would people accept supranational governance?",
     "How do they assess different ways of funding sustainable development?",
     "How would they trade off equality, growth and the environment?"],
    ["Do policymakers endorse the proposals?",
     "What obstacles would the proposals face?",
     "How do policymakers envision sustainable development?",
     "How can research be useful for sustainability policymaking?"],
    ["How can international taxes and transfers transform global distributions and sectors?",
     "How does current inequality translate into varying labour footprints across country-deciles?"],
]
s = prs.slides.add_slide(BLANK); set_bg(s, LIGHT)
slide_title(s, "Research questions")
y0, ch = 1.6, 5.3
for i in range(3):
    wp_column(s, i, y0, ch,
              (lambda col, ix: (lambda x, yb: bullets(s, x + 0.26, yb, COLW - 0.5, ch - 1.3,
                                                       RQ[ix], size=12, bullet_color=col)))(WPS[i]["c"], i))
footer(s, 3)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Methods
# ════════════════════════════════════════════════════════════════════════════════
METHODS = [
    ["Representative survey: 25k respondents across 12 diverse countries (US, JP, RU, SA, CN, IN, BR, NG, EG, DE, FR, UK)",
     "Budget-allocation tasks, Likert scales and open-ended questions",
     "Conjoint analysis on future scenarios",
     "Simulate respondents’ outcomes in the proposals and gauge support"],
    ["Contact ministers’ advisors, diplomats, MPs & MEPs one-by-one to fill a questionnaire and agree to meet",
     "Compare answers from policymakers and the public",
     "Present proposals in interviews and record reactions",
     "Build relationships that help refine the proposals"],
    ["Build on PIK, CIRED and PSE models to represent tax clubs at the country-decile level — combining a CGE, a distributional model and tax & climate modules",
     "This granularity allows computing the labour footprint for any final-demand vector"],
]
s = prs.slides.add_slide(BLANK); set_bg(s, LIGHT)
slide_title(s, "Methods")
for i in range(3):
    wp_column(s, i, y0, ch,
              (lambda col, ix: (lambda x, yb: bullets(s, x + 0.26, yb, COLW - 0.5, ch - 1.3,
                                                       METHODS[ix], size=11.5, bullet_color=col)))(WPS[i]["c"], i))
footer(s, 4)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Impacts
# ════════════════════════════════════════════════════════════════════════════════
IMPACTS = [
    ["Inform society about public attitudes, potentially uncovering support for worldwide game-changing proposals — making them more credible"],
    ["Inform policymakers about the existence and rationale of tax clubs, raising the chance that they are negotiated and implemented",
     "Build fruitful policymaker–researcher relationships, making research questions more policy-relevant"],
    ["Deliver models usable by multilateral organizations to help diplomats negotiate tax clubs",
     "Publish an online labour-footprint simulator that raises awareness of how one’s consumption draws on others’ labour"],
]
s = prs.slides.add_slide(BLANK); set_bg(s, LIGHT)
slide_title(s, "Expected impacts")
for i in range(3):
    wp_column(s, i, y0, ch,
              (lambda col, ix: (lambda x, yb: bullets(s, x + 0.26, yb, COLW - 0.5, ch - 1.3,
                                                       IMPACTS[ix], size=12.5, bullet_color=col)))(WPS[i]["c"], i))
footer(s, 5)

# ════════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Closing
# ════════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); set_bg(s, NAVY)
for i, wp in enumerate(WPS):
    box(s, MSO_SHAPE.RECTANGLE, i * (SW / 3), 0, SW / 3, 0.22, fill=wp["c"])
    box(s, MSO_SHAPE.RECTANGLE, i * (SW / 3), SH - 0.22, SW / 3, 0.22, fill=wp["c"])
text(s, 1.0, 2.55, 11.3, 1.6,
     [[R("Making global redistribution", 34, WHITE, bold=True, font=TITLE_FONT)],
      [R("credible, negotiable and modelled.", 34, WHITE, bold=True, font=TITLE_FONT)]],
     align=PP_ALIGN.CENTER, line_spacing=1.05)
text(s, 1.0, 4.5, 11.3, 0.5,
     [[R("Adrien Fabre", 16, WHITE, bold=True, font=HEAD_FONT),
       R("   ·   CNRS–CIRED   ·   INCENTIVE Horizons", 16, RGBColor(0xB9, 0xC6, 0xD2))]],
     align=PP_ALIGN.CENTER)

prs.save("prez_horizons.pptx")
print("saved prez_horizons.pptx with", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
